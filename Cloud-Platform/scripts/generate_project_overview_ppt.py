from pathlib import Path
from tempfile import TemporaryDirectory

import cairosvg
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "robot-cloud-system-overview.pptx"
ARCH_SVG = ROOT / "docs" / "system-architecture.svg"

FONT = "Noto Sans CJK SC"
BG = RGBColor(247, 248, 243)
INK = RGBColor(22, 48, 47)
MUTED = RGBColor(82, 102, 106)
TEAL = RGBColor(45, 75, 79)
GREEN = RGBColor(86, 163, 123)
BLUE = RGBColor(73, 119, 196)
PURPLE = RGBColor(111, 89, 176)
ORANGE = RGBColor(209, 111, 80)
GOLD = RGBColor(213, 166, 41)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(217, 223, 215)


def set_run(run, size=18, bold=False, color=INK):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def set_text(shape, text, size=18, bold=False, color=INK, align=None):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=INK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.margin_left = Inches(0.03)
    box.text_frame.margin_right = Inches(0.03)
    set_text(box, text, size=size, bold=bold, color=color, align=align)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.6, 0.35, 11.6, 0.45, title, size=25, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, 0.62, 0.88, 11.5, 0.35, subtitle, size=11.5, color=MUTED)
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.24), Inches(1.0), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.fill.background()


def add_bg(slide):
    rect = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG
    rect.line.fill.background()
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)


def add_card(slide, x, y, w, h, title, lines, color=TEAL):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LINE
    card.line.width = Pt(1)
    add_textbox(slide, x + 0.18, y + 0.16, w - 0.36, 0.28, title, size=14, bold=True, color=color)
    ty = y + 0.56
    for line in lines:
        add_textbox(slide, x + 0.22, ty, w - 0.36, 0.28, f"• {line}", size=10.5, color=MUTED)
        ty += 0.33
    return card


def add_pill(slide, x, y, w, text, color):
    pill = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    set_text(pill, text, size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def add_footer(slide, text="Robot Cloud System | 温室表型机器人云端管理平台"):
    add_textbox(slide, 0.6, 7.08, 8.8, 0.22, text, size=8.5, color=RGBColor(110, 126, 128))
    add_textbox(slide, 11.25, 7.08, 1.5, 0.22, "2026-04-24", size=8.5, color=RGBColor(110, 126, 128), align=PP_ALIGN.RIGHT)


def add_step(slide, x, y, w, h, num, title, detail, color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.13), Inches(y + 0.16), Inches(0.34), Inches(0.34)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    set_text(badge, str(num), size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.55, y + 0.12, w - 0.7, 0.28, title, size=12.5, bold=True, color=INK)
    add_textbox(slide, x + 0.18, y + 0.53, w - 0.36, h - 0.62, detail, size=9.3, color=MUTED)


def arrow_between(slide, x, y, w=0.36, color=TEAL):
    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.22)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    with TemporaryDirectory() as td:
        arch_png = Path(td) / "system-architecture.png"
        cairosvg.svg2png(url=str(ARCH_SVG), write_to=str(arch_png), output_width=1800)

        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_textbox(slide, 0.72, 0.72, 8.0, 0.65, "Robot Cloud System", size=35, bold=True, color=INK)
        add_textbox(slide, 0.76, 1.42, 8.6, 0.38, "温室表型机器人云端管理平台项目系统简要讲解", size=17, color=MUTED)
        add_textbox(
            slide,
            0.78,
            2.28,
            5.6,
            0.95,
            "覆盖任务调度、设备控制、采集资产管理、分析结果查询与演示验收的一体化平台。",
            size=20,
            bold=True,
            color=TEAL,
        )
        for i, (text, color) in enumerate(
            [
                ("任务调度", GREEN),
                ("机器人控制", BLUE),
                ("资产管理", ORANGE),
                ("分析结果", PURPLE),
            ]
        ):
            add_pill(slide, 0.8 + i * 1.35, 3.6, 1.05, text, color)
        add_card(
            slide,
            7.45,
            1.22,
            4.85,
            4.95,
            "一句话定位",
            [
                "前端控制台负责操作入口和实时态势展示",
                "Flask API 承载认证、任务、设备、资产和系统接口",
                "MQTT 连接云端与机器人，Celery 执行异步分析",
                "Docker Compose 提供可复现的演示和验收环境",
            ],
            color=GREEN,
        )
        add_footer(slide, "资料来源：README.md、PROJECT_MANUAL.md、docs/SYSTEM_OVERVIEW.md")

        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "1. 项目定位与交付范围", "从温室表型机器人作业出发，形成云端调度、采集、分析、运维闭环")
        add_card(slide, 0.72, 1.7, 3.72, 3.95, "解决的问题", ["机器人任务统一创建与下发", "任务进度、心跳、异常可实时回传", "采集文件与分析结果可追踪", "演示环境可冷启动验收"], color=GREEN)
        add_card(slide, 4.83, 1.7, 3.72, 3.95, "核心能力", ["登录与管理员初始化", "任务、设备、图库、结果、系统管理", "MQTT 任务和命令通道", "上传会话、对象存储、异步分析"], color=BLUE)
        add_card(slide, 8.94, 1.7, 3.72, 3.95, "交付形态", ["Docker Compose 默认演示栈", "环境变量覆盖生产配置", "内置 simulator 完成业务闭环", "smoke / cold-smoke 自动验收"], color=ORANGE)
        add_footer(slide)

        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "2. 总体架构", "浏览器和验收脚本经 Nginx 进入平台，后端连接数据库、对象存储、队列、MQTT 与机器人")
        slide.shapes.add_picture(str(arch_png), Inches(0.52), Inches(1.42), width=Inches(12.2), height=Inches(5.28))
        add_footer(slide)

        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "3. 模块职责", "按入口、应用层、基础设施、设备侧划分边界，降低联调复杂度")
        add_card(slide, 0.65, 1.58, 2.85, 4.75, "Vue 控制台", ["页面：Dashboard、Tasks、Robots、Data、Results、Admin", "通过 /api 调用 REST 接口", "通过 /ws/events 接收事件流"], color=BLUE)
        add_card(slide, 3.7, 1.58, 2.85, 4.75, "Flask API / 服务层", ["统一鉴权和响应格式", "封装任务状态机、机器人命令、上传会话", "生成实时事件和系统告警"], color=PURPLE)
        add_card(slide, 6.75, 1.58, 2.85, 4.75, "基础设施", ["PostgreSQL 保存业务元数据", "Redis 支撑 Celery broker/backend", "MinIO 存资产和结果文件", "Mosquitto 承载 MQTT 通道"], color=GREEN)
        add_card(slide, 9.8, 1.58, 2.85, 4.75, "Worker / Simulator", ["Worker 执行 demo/http/disabled 分析", "定时扫描机器人离线状态", "Simulator 初始化演示管理员", "模拟 ACK、进度、上传和结果闭环"], color=ORANGE)
        add_footer(slide)

        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "4. 核心业务闭环", "从操作员创建任务到前端看到结果，系统依赖 API、MQTT、对象存储和异步队列协作")
        steps = [
            ("登录/初始化", "管理员登录；冷启动时可通过 BOOTSTRAP_TOKEN 初始化首个管理员", GREEN),
            ("创建任务", "API 写入任务、绑定机器人，并进入 PENDING_DISPATCH", BLUE),
            ("MQTT 下发", "向 greenhouse/tasks/<robot>/dispatch 发布任务载荷", PURPLE),
            ("机器人回传", "ACK、心跳、进度、异常和命令事件回到后端", ORANGE),
            ("上传资产", "上传会话写入 MinIO，资产元数据写入 PostgreSQL", GREEN),
            ("异步分析", "Celery worker 执行分析 provider 并写入结果", BLUE),
            ("实时刷新", "WebSocket 推送任务、告警、结果状态到前端", PURPLE),
        ]
        xs = [0.62, 2.42, 4.22, 6.02, 7.82, 9.62, 11.42]
        for i, (title, detail, color) in enumerate(steps):
            add_step(slide, xs[i], 2.05, 1.42, 2.85, i + 1, title, detail, color)
            if i < len(steps) - 1:
                arrow_between(slide, xs[i] + 1.48, 3.35, 0.28, TEAL)
        add_textbox(slide, 0.8, 5.58, 11.8, 0.5, "闭环验收脚本会真实调用 API，等待 MQTT 回执、资产上传、分析结果和告警读取，验证链路完整性。", size=15, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        add_footer(slide)

        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "5. 任务生命周期", "状态机约束任务主路径和异常路径，避免前后端、机器人、分析服务各自解释状态")
        states = [
            "DRAFT",
            "PENDING_DISPATCH",
            "DISPATCHED",
            "ROBOT_ACKED",
            "RUNNING",
            "DATA_UPLOADING",
            "DATA_READY",
            "ANALYZING",
            "COMPLETED",
        ]
        y = 2.05
        for i, state in enumerate(states):
            x = 0.66 + (i % 5) * 2.42
            row_y = y + (i // 5) * 1.2
            color = GREEN if i in (0, 8) else BLUE if i < 5 else PURPLE
            add_pill(slide, x, row_y, 1.72, state, color)
            if i not in (4, 8):
                arrow_between(slide, x + 1.76, row_y + 0.07, 0.34, TEAL)
        add_card(slide, 0.72, 4.88, 3.55, 1.55, "失败与重试", ["FAILED 可回到 PENDING_DISPATCH 或 DATA_READY", "支持按失败阶段重新下发或重新分析"], color=ORANGE)
        add_card(slide, 4.86, 4.88, 3.55, 1.55, "取消路径", ["DISPATCHED / RUNNING / ANALYZING 可进入 CANCELLING", "最终落到 CANCELLED 或 FAILED"], color=GOLD)
        add_card(slide, 9.0, 4.88, 3.55, 1.55, "事件记录", ["TaskEvent、RealtimeEvent、SystemAlert 留痕", "前端通过事件流刷新展示"], color=PURPLE)
        add_footer(slide)

        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "6. 数据与接口设计", "业务数据、文件对象、实时事件分层保存，接口围绕操作闭环组织")
        add_card(slide, 0.7, 1.65, 3.8, 4.75, "核心数据表", ["User / Role：账号与权限", "Robot / RobotHeartbeat / RobotCommand：设备与命令", "Task / TaskEvent：任务和过程事件", "UploadSession / DataAsset：上传会话与资产", "AnalysisJob / AnalysisResult：分析任务和结果"], color=GREEN)
        add_card(slide, 4.78, 1.65, 3.8, 4.75, "关键接口", ["Auth：login、bootstrap-admin、me", "System：health、bootstrap-check、release-readiness", "Tasks：create、dispatch、retry、cancel", "Assets：upload-sessions、complete、query", "Results / Dashboard / Admin：查询和治理"], color=BLUE)
        add_card(slide, 8.86, 1.65, 3.8, 4.75, "通信协议", ["REST API 使用统一响应结构", "WebSocket：/ws/events，首条消息携带 token 认证", "MQTT topic：greenhouse/robots/+/heartbeat 等", "任务下发和机器人命令都走 MQTT publish"], color=ORANGE)
        add_footer(slide)

        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "7. 部署、运行与验收", "默认演示栈可直接启动；生产环境通过 .env 覆盖密钥、存储、MQTT 和分析服务")
        add_card(slide, 0.72, 1.58, 3.72, 4.75, "默认服务", ["backend、worker、frontend、nginx", "postgres、redis、minio、mosquitto", "minio-init 初始化对象桶", "simulator 默认启用"], color=GREEN)
        add_card(slide, 4.83, 1.58, 3.72, 4.75, "常用命令", ["./scripts/docker_stack.sh up", "./scripts/docker_stack.sh smoke", "PROJECT_NAME=robot-cloud-cold ./scripts/docker_stack.sh cold-smoke", "./scripts/docker_stack.sh demo-reset-admin"], color=BLUE)
        add_card(slide, 8.94, 1.58, 3.72, 4.75, "发布门槛", ["health：进程活性", "bootstrap-check：依赖就绪", "release-readiness：生产参数检查", "接入真实机器人时可 --scale simulator=0"], color=ORANGE)
        add_footer(slide)

        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "8. 推荐讲解顺序", "用于 5 到 8 分钟项目汇报：先讲目标，再讲架构，最后用闭环和验收收束")
        add_step(slide, 0.78, 1.62, 3.55, 1.45, 1, "先讲业务目标", "平台面向温室表型机器人，把任务、设备、数据和分析统一到云端。", GREEN)
        add_step(slide, 4.9, 1.62, 3.55, 1.45, 2, "再讲系统架构", "前端经 Nginx 访问 Flask API；后端对接 Postgres、Redis、MinIO、MQTT 和 Worker。", BLUE)
        add_step(slide, 9.02, 1.62, 3.55, 1.45, 3, "重点讲闭环", "创建任务、MQTT 下发、机器人回传、上传资产、异步分析、WebSocket 刷新。", PURPLE)
        add_step(slide, 2.84, 4.05, 3.55, 1.45, 4, "补充可靠性", "状态机约束生命周期；健康检查、发布门槛和告警帮助定位问题。", ORANGE)
        add_step(slide, 6.96, 4.05, 3.55, 1.45, 5, "收尾讲交付", "Docker Compose 演示栈和 smoke/cold-smoke 让系统可复现、可验收。", GOLD)
        add_footer(slide, "建议演示入口：http://localhost；验收入口：./scripts/docker_stack.sh smoke")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(create_presentation())
